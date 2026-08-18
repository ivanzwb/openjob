from collections import deque
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ASSET_DIR = Path(r"C:\Users\weibzhao\.cursor\projects\c-Projects-openJob\assets")
SOURCE_IMAGES = [
    ASSET_DIR / "openjob-slide-1-positioning.png",
    ASSET_DIR / "openjob-core-feature-map.png",
    ASSET_DIR / "openjob-slide-3-onboarding-flow.png",
]
OUTPUT_IMAGES = [
    ASSET_DIR / "openjob-slide-1-positioning-nobg.png",
    ASSET_DIR / "openjob-core-feature-map-nobg.png",
    ASSET_DIR / "openjob-slide-3-onboarding-flow-nobg.png",
]


def is_background_pixel(pixel: tuple[int, int, int, int]) -> bool:
    r, g, b, a = pixel
    if a == 0:
        return True
    # 只移除真正的外层画布背景：接近白色，或非常浅的蓝色角落装饰。
    # 这里刻意不把普通浅灰/浅蓝都算成背景，避免吃掉流程图卡片和容器底色。
    nearly_white = r >= 248 and g >= 248 and b >= 248
    very_soft_blue = r >= 225 and g >= 238 and b >= 248
    return nearly_white or very_soft_blue


def remove_connected_background(src: Path, dst: Path) -> None:
    image = Image.open(src).convert("RGBA")
    pixels = image.load()
    width, height = image.size
    seen = bytearray(width * height)
    queue: deque[tuple[int, int]] = deque()

    def push(x: int, y: int) -> None:
        idx = y * width + x
        if seen[idx]:
            return
        seen[idx] = 1
        if is_background_pixel(pixels[x, y]):
            queue.append((x, y))

    # 只从四边的极窄边界开始 flood fill。
    # 如果从整条边无差别扩张且阈值太宽，会误进流程图内部的浅色块。
    for x in range(width):
        push(x, 0)
        push(x, height - 1)
    for y in range(height):
        push(0, y)
        push(width - 1, y)

    while queue:
        x, y = queue.popleft()
        r, g, b, _ = pixels[x, y]
        pixels[x, y] = (r, g, b, 0)
        if x > 0:
            push(x - 1, y)
        if x + 1 < width:
            push(x + 1, y)
        if y > 0:
            push(x, y - 1)
        if y + 1 < height:
            push(x, y + 1)

    image.save(dst)


def build_pptx(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for image_path in OUTPUT_IMAGES:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)


if __name__ == "__main__":
    for src, dst in zip(SOURCE_IMAGES, OUTPUT_IMAGES):
        remove_connected_background(src, dst)
    build_pptx("docs/openjob-user-intro-image-slides-nobg.pptx")
