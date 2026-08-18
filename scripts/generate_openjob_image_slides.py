from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


SLIDES = [
    Path(r"C:\Users\weibzhao\.cursor\projects\c-Projects-openJob\assets\openjob-slide-1-positioning.png"),
    Path(r"C:\Users\weibzhao\.cursor\projects\c-Projects-openJob\assets\openjob-core-feature-map.png"),
    Path(r"C:\Users\weibzhao\.cursor\projects\c-Projects-openJob\assets\openjob-slide-3-onboarding-flow.png"),
]


def build_deck(output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    for image_path in SLIDES:
        if not image_path.exists():
            raise FileNotFoundError(image_path)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_path)


if __name__ == "__main__":
    build_deck("docs/openjob-user-intro-image-slides.pptx")
