from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def rgb(hex_color: str) -> RGBColor:
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size=20, bold=False, color="0F172A", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = "Microsoft YaHei"
    return box


def add_card(slide, x, y, w, h, title, body, accent="0EA5E9"):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb("F8FAFC")
    card.line.color.rgb = rgb("E2E8F0")
    card.line.width = Pt(1.2)

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = rgb(accent)
    accent_bar.line.fill.background()

    add_textbox(slide, x + 0.2, y + 0.12, w - 0.3, 0.3, title, size=16, bold=True, color="0F172A")
    add_textbox(slide, x + 0.2, y + 0.45, w - 0.3, h - 0.55, body, size=12, color="334155")


def add_arrow(slide, x1, y1, x2, y2, color="0284C7"):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2)
    line.line.end_arrowhead = True


def build_deck(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: what and why
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("F8FAFC")
    bg.line.fill.background()

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = rgb("0F172A")
    top_band.line.fill.background()

    add_textbox(slide, 0.6, 0.25, 8.8, 0.5, "OpenJob：把面试准备变成可执行的日常流程", size=30, bold=True, color="FFFFFF")
    add_textbox(slide, 0.62, 0.82, 8.8, 0.3, "给工具用户的 3 分钟介绍", size=14, color="CBD5E1")

    # right visual motif
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.9), Inches(0.2), Inches(2.8), Inches(2.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb("0EA5E9")
    circle.fill.transparency = 20
    circle.line.fill.background()
    add_textbox(slide, 10.45, 1.1, 1.9, 0.8, "用户\n价值", size=20, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)

    add_card(
        slide,
        0.7,
        1.7,
        3.9,
        2.1,
        "设计思路 1：流程优先",
        "从“岗位目标 -> 学习计划 -> 练习反馈 -> 话术沉淀”串成一条线，减少碎片化工具切换。",
        accent="0284C7",
    )
    add_card(
        slide,
        4.85,
        1.7,
        3.9,
        2.1,
        "设计思路 2：输出导向",
        "不只收集资料，更强调“能说出来”。每一步都围绕面试可复述内容做积累。",
        accent="0EA5E9",
    )
    add_card(
        slide,
        0.7,
        4.0,
        8.05,
        2.35,
        "设计思路 3：人机协作而非替代",
        "AI 负责生成和压缩信息；用户负责修订、标记、复述。最终沉淀的是你自己的话术，而不是模型原文。",
        accent="14B8A6",
    )

    # Slide 2: feature map
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("FFFFFF")
    bg.line.fill.background()

    add_textbox(slide, 0.6, 0.35, 8.4, 0.6, "核心功能地图（用户视角）", size=28, bold=True, color="0F172A")
    add_textbox(slide, 0.62, 0.9, 8.6, 0.3, "每个模块都服务同一个目标：更快形成可复述的面试表达", size=13, color="64748B")

    add_card(slide, 0.8, 1.5, 2.9, 2.05, "01 简历中心", "导入/编辑简历\n多模板预览与导出\n移动端同步查看", accent="0EA5E9")
    add_card(slide, 4.1, 1.5, 2.9, 2.05, "02 目标岗位", "岗位拆解成知识点\n自动生成学习路径\n进度可视化追踪", accent="14B8A6")
    add_card(slide, 7.4, 1.5, 2.9, 2.05, "03 讲解与标注", "划词高亮/记笔记\n细化讲解与去重\n沉淀重点片段", accent="22C55E")
    add_card(slide, 0.8, 4.0, 2.9, 2.05, "04 真题练习", "考我模式即时反馈\n盲区题专项回练\n输出复盘建议", accent="F59E0B")
    add_card(slide, 4.1, 4.0, 2.9, 2.05, "05 话术库", "一键存入表达片段\n持续改写成个人版本\n支持检索与导出", accent="F97316")
    add_card(slide, 7.4, 4.0, 2.9, 2.05, "06 多端与更新", "桌面 + 移动协同\n同步配对与冲突处理\n版本更新入口清晰", accent="8B5CF6")

    # Slide 3: quick start flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("0F172A")
    bg.line.fill.background()

    add_textbox(slide, 0.6, 0.35, 9.0, 0.6, "新用户 10 分钟上手路径", size=30, bold=True, color="FFFFFF")
    add_textbox(slide, 0.62, 0.9, 9.6, 0.3, "按下面 4 步走完，就能开始一轮完整面试准备闭环", size=13, color="CBD5E1")

    steps = [
        ("1", "导入简历", "上传已有简历\n确认字段与模板"),
        ("2", "建立岗位", "选择目标岗位\n生成学习任务"),
        ("3", "做讲解+练习", "划词标注重点\n做考我并复盘"),
        ("4", "沉淀话术", "把可用表达\n存入话术库"),
    ]
    x_positions = [0.8, 3.9, 7.0, 10.1]
    for i, (num, title, body) in enumerate(steps):
        x = x_positions[i]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.45), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("111827")
        card.line.color.rgb = rgb("334155")

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.9), Inches(2.2), Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb("0EA5E9")
        badge.line.fill.background()
        add_textbox(slide, x + 1.11, 2.37, 0.25, 0.2, num, size=13, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)

        add_textbox(slide, x + 0.25, 2.95, 2.0, 0.4, title, size=17, bold=True, color="F8FAFC", align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.25, 3.45, 2.0, 1.5, body, size=12, color="CBD5E1", align=PP_ALIGN.CENTER)

        if i < 3:
            add_arrow(slide, x + 2.5, 3.6, x + 3.0, 3.6, color="38BDF8")

    add_textbox(
        slide,
        0.8,
        5.7,
        11.8,
        1.1,
        "结论：OpenJob 的价值不在“多一个工具”，而在于把准备动作持续转成可复用的话术资产。",
        size=15,
        bold=True,
        color="E2E8F0",
        align=PP_ALIGN.CENTER,
    )

    prs.save(output_path)


if __name__ == "__main__":
    build_deck("docs/openjob-tool-user-intro.pptx")
