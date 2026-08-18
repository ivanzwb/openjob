from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


def rgb(color: str) -> RGBColor:
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


W, H = 13.333333, 7.5
BG = "F7F8FA"
INK = "111827"
MUTED = "5B6472"
LINE = "D9DEE8"
CARD = "FFFFFF"
BLUE = "2F6BFF"
TEAL = "00A6A6"
SOFT_BLUE = "EEF4FF"


def add_text(slide, text, x, y, w, h, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.bold = bold
    return shape


def add_title(slide, title, subtitle):
    add_text(slide, title, 0.65, 0.45, 8.8, 0.48, size=26, bold=True)
    add_text(slide, subtitle, 0.66, 0.98, 8.5, 0.28, size=11.5, color=MUTED)
    add_text(slide, "OpenJob", 11.35, 0.52, 1.3, 0.28, size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, accent=BLUE, num=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = rgb(CARD)
    rect.line.color.rgb = rgb(LINE)
    rect.line.width = Pt(0.8)
    if num:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(y + 0.22), Inches(0.52), Inches(0.32))
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb(SOFT_BLUE)
        pill.line.fill.background()
        add_text(slide, num, x + 0.36, y + 0.27, 0.22, 0.12, size=9, color=accent, bold=True, align=PP_ALIGN.CENTER)
        tx = x + 0.86
        tw = w - 1.05
    else:
        tx = x + 0.28
        tw = w - 0.5
    add_text(slide, title, tx, y + 0.22, tw, 0.28, size=14.5, bold=True)
    add_text(slide, body, x + 0.28, y + 0.68, w - 0.5, h - 0.85, size=10.7, color=MUTED)


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=1.3):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True


def make_deck(path: str):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BG)
    add_title(s, "OpenJob 是什么", "一个把简历、岗位、学习、练习和话术沉淀串起来的面试准备工作台")
    add_text(s, "从“我该准备什么”到“我能不能说出来”", 0.8, 1.75, 4.9, 0.45, size=19, bold=True)
    add_text(s, "OpenJob 的重点不是替用户准备答案，而是把准备过程拆成可执行动作，并持续沉淀成自己的表达资产。", 0.82, 2.25, 4.75, 0.9, size=12.2, color=MUTED)
    add_card(s, 6.35, 1.62, 2.0, 1.35, "流程优先", "岗位目标\n知识点\n学习任务\n复盘", BLUE, "01")
    add_card(s, 8.65, 1.62, 2.0, 1.35, "输出导向", "每一步最终\n都回到可复述表达", TEAL, "02")
    add_card(s, 10.95, 1.62, 1.75, 1.35, "人机协作", "AI 生成初稿\n用户改成自己的话", BLUE, "03")
    steps = [("简历", 0.95), ("岗位", 3.0), ("学习", 5.05), ("练习", 7.1), ("话术", 9.15)]
    y = 4.55
    for idx, (label, x) in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.35), Inches(0.68))
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(CARD)
        box.line.color.rgb = rgb(LINE)
        add_text(s, label, x + 0.33, y + 0.21, 0.7, 0.2, size=13, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_arrow(s, x + 1.43, y + 0.34, steps[idx + 1][1] - 0.1, y + 0.34, BLUE)
    add_text(s, "准备闭环", 10.9, 4.78, 1.6, 0.3, size=12, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

    # Slide 2
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BG)
    add_title(s, "核心功能地图", "按用户完成面试准备的顺序组织，而不是按技术模块堆叠")
    add_card(s, 0.8, 1.6, 3.65, 1.25, "简历中心", "编辑经历、选择模板、预览并导出 PDF", BLUE, "01")
    add_card(s, 4.85, 1.6, 3.65, 1.25, "目标岗位与学习计划", "输入岗位信息，拆成知识点和学习任务", BLUE, "02")
    add_card(s, 8.9, 1.6, 3.65, 1.25, "讲解与标注", "看讲解、划重点、记笔记、细化难点", BLUE, "03")
    add_card(s, 0.8, 3.45, 3.65, 1.25, "练习与复盘", "考我、模拟面试、定位盲区并回练", TEAL, "04")
    add_card(s, 4.85, 3.45, 3.65, 1.25, "话术库", "把可用表达收集、检索、复用和导出", TEAL, "05")
    add_card(s, 8.9, 3.45, 3.65, 1.25, "多端与同步", "桌面端深度编辑，移动端随时复习，同步保持一致", TEAL, "06")
    for x1, y1, x2, y2 in [(4.45, 2.22, 4.78, 2.22), (8.5, 2.22, 8.83, 2.22), (4.45, 4.07, 4.78, 4.07)]:
        add_arrow(s, x1, y1, x2, y2, BLUE, 1.1)
    add_text(s, "设计判断：功能不是孤立入口，而是围绕“形成自己的面试表达”形成连续路径。", 0.95, 5.65, 11.2, 0.42, size=14, color=INK, bold=True, align=PP_ALIGN.CENTER)

    # Slide 3
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BG)
    add_title(s, "新用户 10 分钟上手路径", "第一次使用时，不需要先理解所有功能，按 4 步走完即可进入准备闭环")
    flow = [
        ("导入或创建简历", "补全经历\n确认模板和导出效果"),
        ("创建目标岗位", "输入公司 / 岗位 / JD\n生成知识点与任务"),
        ("学习讲解并做标注", "看讲解，划重点\n记笔记，细化难点"),
        ("练习并沉淀话术", "用考我暴露盲区\n把好表达存入话术库"),
    ]
    xs = [0.8, 3.9, 7.0, 10.1]
    for i, (title, body) in enumerate(flow):
        add_card(s, xs[i], 2.0, 2.35, 2.3, title, body, BLUE if i < 2 else TEAL, f"0{i+1}")
        if i < 3:
            add_arrow(s, xs[i] + 2.42, 3.15, xs[i + 1] - 0.12, 3.15, BLUE)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.35), Inches(5.45), Inches(10.6), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb("FFFFFF")
    bar.line.color.rgb = rgb(LINE)
    add_text(s, "建议节奏：学 1 个点  ->  练 1 次  ->  存 1 条话术", 2.05, 5.68, 9.2, 0.2, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)

    prs.save(path)


if __name__ == "__main__":
    make_deck("docs/openjob-user-intro-clean-slides.pptx")
